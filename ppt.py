import collections 
import collections.abc
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Pt


prs = Presentation('./ppt-templates/explodingkittens.pptx')
most_recent_not_shape = 0

text_runs = []
""" for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            if most_recent_not_shape != 0:
                if most_recent_not_shape.name == "Header Textbox":
                    for nested_shape in most_recent_not_shape.shapes:
                        if nested_shape.has_text_frame:
                           print(dir(nested_shape))


            most_recent_not_shape = shape
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                text_runs.append(run.text)
                if(run.text == 'Additional text or pictures go here'):
                    print('test')
                print(run.text)
prs.save('test.pptx') """

slides = prs.slides

counter = 0
allSlides = []
for slide in prs.slides:
    questionAndAnswer = []
    for shape in slide.shapes:
        if shape.name == "Header Textbox" or shape.name == "Answer Textbox":
            for nested_shape in shape.shapes:
                if nested_shape.has_text_frame:
                    default = nested_shape.text_frame.paragraphs[0].runs[0]       
                    questionAndAnswer.append(default)
        if shape.name == "Answer Textbox":
            allSlides.append(questionAndAnswer)
for q, a in allSlides:
    print(q.text)
    print(a.text)
prs.save('game-test.pptx')
